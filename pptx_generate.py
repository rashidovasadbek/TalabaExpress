"""
Zamonaviy PPTX generatori — SHABLON + RANG tizimi.

Ikki bosqichli tanlov:
  1) TEMPLATE (shablon uslubi) — 10 xil: butun prezentatsiyaning umumiy ko'rinishi
     (fon, shrift, layout oilasi, bullet uslubi, bezak).
  2) THEME (rang/dizayn) — 15 xil rang palitrasi.

Har bir slayd shu shablon uslubida IZCHIL chiziladi (tasodifiy emas).
Tashqi .pptx shablonga bog'liq emas — barchasi qo'lda chiziladi.
"""

import os
import uuid
import random

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.chart.data import CategoryChartData

from ai_service import GeminiService

try:
    from PIL import Image
    _HAS_PIL = True
except Exception:
    _HAS_PIL = False

WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CREAM = RGBColor(0xF6, 0xF1, 0xE7)
LIGHT_MUTED = RGBColor(0xAA, 0xB2, 0xC0)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ---------------------------------------------------------------------------
# 15 ta RANG palitrasi
# ---------------------------------------------------------------------------
def _C(r, g, b):
    return RGBColor(r, g, b)


def _mk(primary, accent, accent2, text, light, name):
    return {"name": name, "primary": _C(*primary), "accent": _C(*accent),
            "accent2": _C(*accent2), "text": _C(*text), "muted": _C(0x6B, 0x72, 0x80),
            "light": _C(*light), "on_dark": WHITE}


THEMES = {
    "ocean":    _mk((0x12, 0x2A, 0x4A), (0x3B, 0x82, 0xF6), (0x60, 0xA5, 0xFA), (0x1F, 0x2A, 0x37), (0xEF, 0xF4, 0xFB), "Okean"),
    "emerald":  _mk((0x0B, 0x3D, 0x2E), (0x10, 0xB9, 0x81), (0x34, 0xD3, 0x99), (0x18, 0x2A, 0x24), (0xEC, 0xFD, 0xF5), "Zumrad"),
    "sunset":   _mk((0x3B, 0x12, 0x3A), (0xF9, 0x73, 0x16), (0xFB, 0x92, 0x3C), (0x2A, 0x18, 0x29), (0xFD, 0xF2, 0xF8), "Shafaq"),
    "crimson":  _mk((0x4A, 0x0E, 0x12), (0xE1, 0x1D, 0x48), (0xFB, 0x71, 0x85), (0x2A, 0x14, 0x17), (0xFF, 0xF1, 0xF2), "Qirmizi"),
    "violet":   _mk((0x2E, 0x10, 0x65), (0x8B, 0x5C, 0xF6), (0xA7, 0x8B, 0xFA), (0x1E, 0x15, 0x35), (0xF5, 0xF3, 0xFF), "Binafsha"),
    "teal":     _mk((0x0A, 0x3A, 0x3A), (0x14, 0xB8, 0xA6), (0x2D, 0xD4, 0xBF), (0x14, 0x29, 0x2A), (0xF0, 0xFD, 0xFA), "Moviy-yashil"),
    "amber":    _mk((0x42, 0x20, 0x06), (0xF5, 0x9E, 0x0B), (0xFB, 0xBF, 0x24), (0x2A, 0x1C, 0x0A), (0xFF, 0xFB, 0xEB), "Kahrabo"),
    "indigo":   _mk((0x1E, 0x1B, 0x4B), (0x63, 0x66, 0xF1), (0x81, 0x8C, 0xF8), (0x1A, 0x1A, 0x2E), (0xEE, 0xF2, 0xFF), "Indigo"),
    "rose":     _mk((0x4C, 0x05, 0x19), (0xF4, 0x3F, 0x5E), (0xFB, 0x71, 0x85), (0x2A, 0x0E, 0x18), (0xFF, 0xF1, 0xF2), "Pushti"),
    "slate":    _mk((0x1E, 0x29, 0x3B), (0x38, 0xBD, 0xF8), (0x7D, 0xD3, 0xFC), (0x1E, 0x29, 0x3B), (0xF1, 0xF5, 0xF9), "Kulrang-ko'k"),
    "forest":   _mk((0x14, 0x36, 0x1F), (0x22, 0xC5, 0x5E), (0x4A, 0xDE, 0x80), (0x14, 0x27, 0x1A), (0xF0, 0xFD, 0xF4), "O'rmon"),
    "midnight": _mk((0x0B, 0x12, 0x20), (0x38, 0xBD, 0xF8), (0x81, 0x8C, 0xF8), (0x11, 0x18, 0x27), (0xEE, 0xF2, 0xFF), "Yarim tun"),
    "coral":    _mk((0x4A, 0x14, 0x10), (0xF9, 0x73, 0x62), (0xFC, 0xA5, 0xA5), (0x2A, 0x14, 0x10), (0xFF, 0xF1, 0xF0), "Marjon"),
    "cyan":     _mk((0x08, 0x33, 0x44), (0x06, 0xB6, 0xD4), (0x22, 0xD3, 0xEE), (0x0E, 0x2A, 0x33), (0xEC, 0xFE, 0xFF), "Siyohrang"),
    "plum":     _mk((0x3B, 0x07, 0x64), (0xC0, 0x26, 0xD3), (0xE8, 0x79, 0xF9), (0x2A, 0x10, 0x33), (0xFD, 0xF4, 0xFF), "Olxo'ri"),
}


# ---------------------------------------------------------------------------
# 10 ta SHABLON uslubi
#   bg:       light | dark | cream
#   title:    titul slayd varianti
#   content:  kontent layout oilasi
#   bullet:   arrow | dash | number | none
#   font:     shrift nomi
# ---------------------------------------------------------------------------
TEMPLATES = {
    "classic":     {"name": "Klassik",     "bg": "light", "title": "block",      "content": "classic",  "bullet": "dash",   "font": "Georgia"},
    "minimalist":  {"name": "Minimalist",  "bg": "light", "title": "minimal",    "content": "minimal",  "bullet": "none",   "font": "Calibri"},
    "bold":        {"name": "Bold",        "bg": "light", "title": "split",      "content": "split",    "bullet": "arrow",  "font": "Arial"},
    "corporate":   {"name": "Korporativ",  "bg": "light", "title": "band",       "content": "band",     "bullet": "arrow",  "font": "Calibri"},
    "modern":      {"name": "Zamonaviy",   "bg": "light", "title": "geo",        "content": "grid",     "bullet": "card",   "font": "Calibri"},
    "dark":        {"name": "Tungi",       "bg": "dark",  "title": "centerdark", "content": "standard", "bullet": "arrow",  "font": "Calibri"},
    "creative":    {"name": "Ijodiy",      "bg": "light", "title": "diagonal",   "content": "creative", "bullet": "arrow",  "font": "Calibri"},
    "elegant":     {"name": "Nafis",       "bg": "cream", "title": "elegant",    "content": "elegant",  "bullet": "dash",   "font": "Georgia"},
    "infographic": {"name": "Infografik",  "bg": "light", "title": "stat",       "content": "numbered", "bullet": "number", "font": "Calibri"},
    "photo":       {"name": "Rasm asosida","bg": "light", "title": "frame",      "content": "photo",    "bullet": "arrow",  "font": "Calibri"},
}

ICON_MAP = [
    (("kirish", "introduction", "muqaddima"), "📖"),
    (("xulosa", "conclusion", "yakun"), "✅"),
    (("tarix", "history"), "🏛"),
    (("ta'rif", "tushuncha", "concept", "asosiy"), "💡"),
    (("texnologi", "technology", "tizim"), "⚙"),
    (("afzal", "advantage", "foyda"), "⭐"),
    (("muammo", "problem", "kamchilik"), "⚠"),
    (("statistik", "data", "ko'rsatkich", "ma'lumot"), "📊"),
    (("kelajak", "future", "istiqbol"), "🚀"),
    (("ta'lim", "education", "o'qit"), "🎓"),
]

CHART_TYPE_MAP = {
    "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "bar":    XL_CHART_TYPE.BAR_CLUSTERED,
    "pie":    XL_CHART_TYPE.PIE,
    "line":   XL_CHART_TYPE.LINE_MARKERS,
}


# ---------------------------------------------------------------------------
# Render konteksti (rang + shablon birgalikda)
# ---------------------------------------------------------------------------
def _ctx(theme, tpl):
    bgmode = tpl["bg"]
    if bgmode == "dark":
        bg, txt, titlec, footer, line = theme["primary"], theme["on_dark"], theme["on_dark"], LIGHT_MUTED, theme["accent2"]
    elif bgmode == "cream":
        bg, txt, titlec, footer, line = CREAM, theme["text"], theme["primary"], theme["muted"], theme["light"]
    else:
        bg, txt, titlec, footer, line = WHITE, theme["text"], theme["primary"], theme["muted"], theme["light"]
    return {
        "theme": theme, "font": tpl["font"], "bg": bg, "text": txt, "title": titlec,
        "accent": theme["accent"], "accent2": theme["accent2"], "primary": theme["primary"],
        "on_dark": theme["on_dark"], "footer": footer, "line": line,
        "bullet": tpl["bullet"], "family": tpl["content"], "title_kind": tpl["title"],
    }


# ---------------------------------------------------------------------------
# Past darajadagi yordamchilar
# ---------------------------------------------------------------------------
def _blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _fill_background(slide, color):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = color


def _add_rect(slide, l, t, w, h, color, shape=MSO_SHAPE.RECTANGLE):
    sh = slide.shapes.add_shape(shape, l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def _add_text(slide, l, t, w, h, text, size, color, bold=False, align=PP_ALIGN.LEFT,
              anchor=MSO_ANCHOR.TOP, font="Calibri", italic=False, wrap=True):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    f = r.font
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.name = font
    f.color.rgb = color
    return box, tf


def _clean_bullets(raw):
    out = []
    for line in (raw or "").split("\n"):
        x = line.strip()
        if not x:
            continue
        x = x.lstrip("*-•▸●○◦·> ").strip().replace("**", "").strip()
        if x:
            out.append(x)
    return out


def _shorten(text, limit):
    text = (text or "").strip()
    if len(text) <= limit:
        return text
    return text[:limit].rsplit(" ", 1)[0].rstrip(",;:.") + "…"


def _pick_icon(title):
    t = (title or "").lower()
    for keys, emoji in ICON_MAP:
        if any(k in t for k in keys):
            return emoji
    return "📌"


def _marker(ctx, i):
    m = ctx["bullet"]
    if m == "number":
        return f"{i+1}.  "
    if m == "dash":
        return "—  "
    if m == "none":
        return ""
    return "▸  "


def _add_image_fit(slide, path, bl, bt, bw, bh):
    try:
        ratio = None
        if _HAS_PIL:
            with Image.open(path) as im:
                iw, ih = im.size
                if iw and ih:
                    ratio = iw / ih
        if ratio is None:
            ratio = bw / bh
        br = bw / bh
        if ratio > br:
            w = bw; h = int(bw / ratio)
        else:
            h = bh; w = int(bh * ratio)
        slide.shapes.add_picture(path, bl + (bw - w) // 2, bt + (bh - h) // 2, width=w, height=h)
        return True
    except Exception as e:
        print(f"⚠️ fit rasm xato: {e}")
        return False


def _add_image_cover(slide, path, l, t, w, h):
    try:
        if _HAS_PIL:
            with Image.open(path) as im:
                iw, ih = im.size
                target = w / h; src = iw / ih
                if src > target:
                    nw = int(ih * target); x0 = (iw - nw) // 2
                    im = im.crop((x0, 0, x0 + nw, ih))
                else:
                    nh = int(iw / target); y0 = (ih - nh) // 2
                    im = im.crop((0, y0, iw, y0 + nh))
                tmp = path + ".cover.jpg"
                im.convert("RGB").save(tmp)
                slide.shapes.add_picture(tmp, l, t, width=w, height=h)
                return True
        slide.shapes.add_picture(path, l, t, width=w, height=h)
        return True
    except Exception as e:
        print(f"⚠️ cover rasm xato: {e}")
        return _add_image_fit(slide, path, l, t, w, h)


def _img_frame(slide, ctx, path, l, t, w, h):
    _add_rect(slide, l - Inches(0.06), t - Inches(0.06), w + Inches(0.12), h + Inches(0.12), ctx["line"])
    _add_image_cover(slide, path, l, t, w, h)


def _bullets(slide, ctx, items, l, t, w, h, base, align=PP_ALIGN.LEFT, force_marker=None):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    for i, line in enumerate(items):
        line = _shorten(line, 200)
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(9)
        p.line_spacing = 1.08
        mk = force_marker if force_marker is not None else _marker(ctx, i)
        if mk:
            mr = p.add_run()
            mr.text = mk
            mr.font.size = Pt(base)
            mr.font.bold = True
            mr.font.name = ctx["font"]
            mr.font.color.rgb = ctx["accent"]
        r = p.add_run()
        r.text = line
        r.font.size = Pt(base)
        r.font.name = ctx["font"]
        r.font.color.rgb = ctx["text"]
    return box


def _cards(slide, ctx, items, l, t, w, total_h, base, numbered=False):
    n = max(1, len(items))
    gap = Inches(0.14)
    ch = int((total_h - gap * (n - 1)) / n)
    for i, line in enumerate(items):
        y = t + i * (ch + gap)
        _add_rect(slide, l, y, w, ch, ctx["theme"]["light"], shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        _add_rect(slide, l, y, Inches(0.16), ch, ctx["accent"])
        tx = l + Inches(0.4)
        if numbered:
            _add_text(slide, l + Inches(0.3), y, Inches(0.7), ch, str(i + 1),
                      base + 4, ctx["accent"], bold=True, anchor=MSO_ANCHOR.MIDDLE, font=ctx["font"])
            tx = l + Inches(1.0)
        tb = slide.shapes.add_textbox(tx, y, l + w - tx - Inches(0.2), ch)
        tf = tb.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = _shorten(line, 130)
        r.font.size = Pt(base)
        r.font.name = ctx["font"]
        r.font.color.rgb = ctx["theme"]["text"]


def _add_footer(slide, ctx, index, total, topic, show_topic=True):
    _add_rect(slide, Inches(0.9), Inches(6.95), Inches(11.53), Inches(0.02), ctx["line"])
    if show_topic:
        _add_text(slide, Inches(0.9), Inches(7.0), Inches(9.0), Inches(0.4),
                  (topic or "")[:60], 10, ctx["footer"], font=ctx["font"])
    _add_text(slide, Inches(11.0), Inches(7.0), Inches(1.43), Inches(0.4),
              f"{index} / {total}", 10, ctx["footer"], align=PP_ALIGN.RIGHT, font=ctx["font"])


def _header(slide, ctx, title, icons=False):
    """Standart yuqori sarlavha (aksent kvadrat + sarlavha + nozik chiziq)."""
    _add_rect(slide, Inches(0.9), Inches(0.62), Inches(0.18), Inches(0.7), ctx["accent"])
    tl = Inches(1.25)
    if icons:
        _add_text(slide, Inches(1.25), Inches(0.45), Inches(0.9), Inches(1.0),
                  _pick_icon(title), 26, ctx["title"], anchor=MSO_ANCHOR.MIDDLE, font="Segoe UI Emoji")
        tl = Inches(2.1)
    _add_text(slide, tl, Inches(0.5), Inches(10.5), Inches(1.0),
              title, 27, ctx["title"], bold=True, anchor=MSO_ANCHOR.MIDDLE, font=ctx["font"])
    _add_rect(slide, Inches(0.9), Inches(1.55), Inches(11.53), Inches(0.03), ctx["line"])


# ---------------------------------------------------------------------------
# Kontent layout OILALARI  (s, ctx, title, bullets, hi, ip, icons, parity)
# ---------------------------------------------------------------------------
def _fam_standard(s, ctx, title, b, hi, ip, icons, parity):
    _header(s, ctx, title, icons=icons)
    if hi and parity % 2 == 0:
        _bullets(s, ctx, b[:4], Inches(1.0), Inches(1.95), Inches(6.4), Inches(4.8), 16)
        _img_frame(s, ctx, ip, Inches(7.75), Inches(1.95), Inches(4.6), Inches(4.5))
    elif hi:
        _img_frame(s, ctx, ip, Inches(0.95), Inches(1.95), Inches(4.6), Inches(4.5))
        _bullets(s, ctx, b[:4], Inches(5.95), Inches(1.95), Inches(6.4), Inches(4.8), 16)
    else:
        _bullets(s, ctx, b[:6], Inches(1.0), Inches(1.95), Inches(11.3), Inches(4.8), 17)


def _fam_band(s, ctx, title, b, hi, ip, icons, parity):
    _add_rect(s, Inches(0), Inches(0), SLIDE_W, Inches(1.5), ctx["primary"])
    _add_rect(s, Inches(0), Inches(1.5), SLIDE_W, Inches(0.08), ctx["accent"])
    tl = Inches(0.9)
    if icons:
        _add_text(s, Inches(0.9), Inches(0.25), Inches(0.9), Inches(1.0), _pick_icon(title),
                  26, ctx["on_dark"], anchor=MSO_ANCHOR.MIDDLE, font="Segoe UI Emoji")
        tl = Inches(1.7)
    _add_text(s, tl, Inches(0.3), Inches(11.0), Inches(0.95), title, 26, ctx["on_dark"],
              bold=True, anchor=MSO_ANCHOR.MIDDLE, font=ctx["font"])
    if hi:
        _bullets(s, ctx, b[:4], Inches(1.0), Inches(1.95), Inches(6.4), Inches(4.7), 16)
        _img_frame(s, ctx, ip, Inches(7.75), Inches(1.95), Inches(4.6), Inches(4.4))
    else:
        _bullets(s, ctx, b[:6], Inches(1.0), Inches(1.95), Inches(11.3), Inches(4.7), 17)


def _fam_sidepanel(s, ctx, title, b, hi, ip, icons, parity):
    _add_rect(s, Inches(0), Inches(0), Inches(4.3), SLIDE_H, ctx["primary"])
    _add_rect(s, Inches(4.3), Inches(0), Inches(0.08), SLIDE_H, ctx["accent"])
    if icons:
        _add_text(s, Inches(0.6), Inches(0.7), Inches(1.0), Inches(1.0), _pick_icon(title),
                  30, ctx["on_dark"], anchor=MSO_ANCHOR.MIDDLE, font="Segoe UI Emoji")
    _add_text(s, Inches(0.6), Inches(1.9), Inches(3.3), Inches(3.0), title, 26, ctx["on_dark"], bold=True, font=ctx["font"])
    _add_rect(s, Inches(0.62), Inches(5.0), Inches(1.6), Inches(0.08), ctx["accent2"])
    if hi:
        _add_image_cover(s, ip, Inches(0.6), Inches(5.3), Inches(3.1), Inches(1.6))
    _bullets(s, ctx, b[:6], Inches(4.8), Inches(1.6), Inches(7.7), Inches(4.9), 17)


def _fam_split(s, ctx, title, b, hi, ip, icons, parity):
    _add_rect(s, Inches(0), Inches(0), Inches(5.4), SLIDE_H, ctx["primary"])
    _add_text(s, Inches(0.7), Inches(1.4), Inches(4.2), Inches(2.8), title, 32, ctx["on_dark"], bold=True, font=ctx["font"])
    _add_rect(s, Inches(0.72), Inches(4.4), Inches(1.8), Inches(0.1), ctx["accent"])
    if hi:
        _add_image_cover(s, ip, Inches(0.7), Inches(4.8), Inches(4.0), Inches(2.0))
    _bullets(s, ctx, b[:5], Inches(5.9), Inches(1.6), Inches(6.6), Inches(4.9), 17)


def _fam_cards(s, ctx, title, b, hi, ip, icons, parity):
    _header(s, ctx, title, icons=icons)
    numbered = ctx["bullet"] == "number"
    if hi:
        _cards(s, ctx, b[:4], Inches(1.0), Inches(1.95), Inches(6.4), Inches(4.6), 15, numbered)
        _img_frame(s, ctx, ip, Inches(7.75), Inches(1.95), Inches(4.6), Inches(4.5))
    else:
        _cards(s, ctx, b[:5], Inches(1.0), Inches(1.95), Inches(11.3), Inches(4.7), 16, numbered)


def _fam_numbered(s, ctx, title, b, hi, ip, icons, parity):
    """Infografik — raqamli doira + matn qatorlari (+ rasm)."""
    _header(s, ctx, title, icons=icons)
    items = b[:4] if hi else b[:5]
    n = max(1, len(items))
    area_w = Inches(6.4) if hi else Inches(11.3)
    top = Inches(2.0)
    gap = Inches(0.25)
    rh = (Inches(4.6) - gap * (n - 1)) / n
    for i, line in enumerate(items):
        y = top + i * (rh + gap)
        _add_rect(s, Inches(1.0), y + rh / 2 - Inches(0.32), Inches(0.64), Inches(0.64), ctx["accent"], shape=MSO_SHAPE.OVAL)
        _add_text(s, Inches(1.0), y + rh / 2 - Inches(0.32), Inches(0.64), Inches(0.64), str(i + 1), 18, ctx["on_dark"], bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=ctx["font"])
        tb = s.shapes.add_textbox(Inches(1.85), y, area_w - Inches(0.9), rh)
        tf = tb.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = _shorten(line, 120)
        r.font.size = Pt(15)
        r.font.name = ctx["font"]
        r.font.color.rgb = ctx["text"]
    if hi:
        _img_frame(s, ctx, ip, Inches(7.75), Inches(1.95), Inches(4.6), Inches(4.5))


def _fam_classic(s, ctx, title, b, hi, ip, icons, parity):
    """Klassik — sarlavha + bulletlar chapda, o'rtada nozik ajratuvchi, rasm o'ngda."""
    _header(s, ctx, title, icons=icons)
    if hi:
        _bullets(s, ctx, b[:4], Inches(1.0), Inches(1.95), Inches(6.0), Inches(4.8), 16)
        _add_rect(s, Inches(7.25), Inches(2.1), Inches(0.02), Inches(4.4), ctx["line"])
        _img_frame(s, ctx, ip, Inches(7.7), Inches(1.95), Inches(4.65), Inches(4.5))
    else:
        _bullets(s, ctx, b[:6], Inches(1.0), Inches(1.95), Inches(11.3), Inches(4.8), 17)


def _fam_grid(s, ctx, title, b, hi, ip, icons, parity):
    """Zamonaviy — 2×2 karta to'ri."""
    _header(s, ctx, title, icons=icons)
    items = b[:4]
    while len(items) < 1:
        items.append("")
    cw = Inches(5.55)
    chh = Inches(2.25)
    gx = Inches(0.25)
    gy = Inches(0.25)
    x0 = Inches(1.0)
    y0 = Inches(2.0)
    for i, line in enumerate(items[:4]):
        col = i % 2
        row = i // 2
        x = x0 + col * (cw + gx)
        y = y0 + row * (chh + gy)
        _add_rect(s, x, y, cw, chh, ctx["theme"]["light"], shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        _add_rect(s, x, y, cw, Inches(0.14), ctx["accent"])
        tb = s.shapes.add_textbox(x + Inches(0.3), y + Inches(0.3), cw - Inches(0.6), chh - Inches(0.5))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = _shorten(line, 110)
        r.font.size = Pt(15)
        r.font.name = ctx["font"]
        r.font.color.rgb = ctx["theme"]["text"]


def _fam_elegant(s, ctx, title, b, hi, ip, icons, parity):
    """Nafis — markazlashgan sarlavha, markazlashgan dash bulletlar, ramkali rasm."""
    _add_text(s, Inches(1.0), Inches(0.8), Inches(11.3), Inches(0.9), title, 28, ctx["title"], bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=ctx["font"])
    _add_rect(s, Inches(5.66), Inches(1.7), Inches(2.0), Inches(0.03), ctx["accent"])
    if hi:
        _bullets(s, ctx, b[:4], Inches(1.4), Inches(2.2), Inches(6.0), Inches(4.4), 16, force_marker="—  ")
        _img_frame(s, ctx, ip, Inches(7.9), Inches(2.3), Inches(4.3), Inches(4.0))
    else:
        _bullets(s, ctx, b[:5], Inches(2.5), Inches(2.2), Inches(8.3), Inches(4.4), 18, align=PP_ALIGN.CENTER, force_marker="—  ")


def _fam_minimal(s, ctx, title, b, hi, ip, icons, parity):
    _add_rect(s, Inches(1.0), Inches(0.9), Inches(0.5), Inches(0.07), ctx["accent"])
    _add_text(s, Inches(1.0), Inches(1.1), Inches(11.0), Inches(1.0), title, 30, ctx["title"], bold=True, font=ctx["font"])
    if hi and parity % 2 == 0:
        _bullets(s, ctx, b[:4], Inches(1.0), Inches(2.6), Inches(6.6), Inches(4.0), 17)
        _add_image_cover(s, ip, Inches(8.0), Inches(2.6), Inches(4.3), Inches(3.8))
    else:
        _bullets(s, ctx, b[:5], Inches(1.0), Inches(2.6), Inches(11.0), Inches(4.0), 18)


def _fam_creative(s, ctx, title, b, hi, ip, icons, parity):
    _add_rect(s, Inches(0), Inches(0), Inches(1.4), Inches(1.4), ctx["accent"], shape=MSO_SHAPE.RIGHT_TRIANGLE)
    tri = _add_rect(s, Inches(12.0), Inches(6.1), Inches(1.33), Inches(1.4), ctx["accent2"], shape=MSO_SHAPE.RIGHT_TRIANGLE)
    tri.rotation = 180
    _add_text(s, Inches(1.0), Inches(0.7), Inches(11.0), Inches(1.0), title, 28, ctx["title"], bold=True, anchor=MSO_ANCHOR.MIDDLE, font=ctx["font"])
    _add_rect(s, Inches(1.0), Inches(1.7), Inches(2.2), Inches(0.08), ctx["accent"])
    if hi:
        # Rasm CHAPDA (boshqa shablonlardan farqlash uchun), nuqta markerli bulletlar o'ngda
        _add_image_cover(s, ip, Inches(1.0), Inches(2.1), Inches(4.7), Inches(4.4))
        _bullets(s, ctx, b[:4], Inches(6.2), Inches(2.1), Inches(6.2), Inches(4.4), 16, force_marker="●  ")
    else:
        _bullets(s, ctx, b[:6], Inches(1.0), Inches(2.1), Inches(11.3), Inches(4.5), 17, force_marker="●  ")


def _fam_photo(s, ctx, title, b, hi, ip, icons, parity):
    if hi and parity % 2 == 0:
        _add_image_cover(s, ip, Inches(0), Inches(0), Inches(5.2), SLIDE_H)
        _add_text(s, Inches(5.7), Inches(0.7), Inches(7.0), Inches(1.0), title, 27, ctx["title"], bold=True, anchor=MSO_ANCHOR.MIDDLE, font=ctx["font"])
        _add_rect(s, Inches(5.7), Inches(1.7), Inches(2.0), Inches(0.08), ctx["accent"])
        _bullets(s, ctx, b[:5], Inches(5.7), Inches(2.0), Inches(7.0), Inches(4.6), 16)
    elif hi:
        _add_image_cover(s, ip, Inches(8.13), Inches(0), Inches(5.2), SLIDE_H)
        _add_text(s, Inches(0.9), Inches(0.7), Inches(6.8), Inches(1.0), title, 27, ctx["title"], bold=True, anchor=MSO_ANCHOR.MIDDLE, font=ctx["font"])
        _add_rect(s, Inches(0.9), Inches(1.7), Inches(2.0), Inches(0.08), ctx["accent"])
        _bullets(s, ctx, b[:5], Inches(0.9), Inches(2.0), Inches(6.8), Inches(4.6), 16)
    else:
        _fam_band(s, ctx, title, b, hi, ip, icons, parity)


_FAMILIES = {
    "standard": _fam_standard, "band": _fam_band, "sidepanel": _fam_sidepanel,
    "split": _fam_split, "cards": _fam_cards, "numbered": _fam_numbered,
    "minimal": _fam_minimal, "creative": _fam_creative, "photo": _fam_photo,
    "classic": _fam_classic, "grid": _fam_grid, "elegant": _fam_elegant,
}
# Quyuq chap panelli oilalar — footer mavzu matni o'chiriladi
_PANEL_FAMILIES = {"sidepanel", "split"}


# ---------------------------------------------------------------------------
# TITUL slayd variantlari
# ---------------------------------------------------------------------------
def _title_info(slide, ctx, d, color):
    info = [f"Bajardi:  {d.get('student_fio', '')}",
            f"Guruh:  {d.get('student_group', '')}",
            f"Yil:  {d.get('year', '')}"]
    box = slide.shapes.add_textbox(Inches(0.9), Inches(5.6), Inches(7.0), Inches(1.5))
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(info):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = line
        r.font.size = Pt(15)
        r.font.name = ctx["font"]
        r.font.color.rgb = color
        p.space_after = Pt(2)


def _title_block(prs, ctx, d):
    s = _blank_slide(prs)
    _fill_background(s, ctx["primary"])
    _add_rect(s, Inches(0), Inches(0), Inches(0.28), SLIDE_H, ctx["accent"])
    _add_rect(s, Inches(11.2), Inches(0), Inches(2.13), Inches(2.2), ctx["accent"])
    _add_rect(s, Inches(10.4), Inches(0), Inches(0.7), Inches(1.3), ctx["accent2"])
    uni = d.get("university_name") or d.get("uni_kafedra", "")
    kaf = d.get("uni_kafedra", "")
    _add_text(s, Inches(0.9), Inches(0.9), Inches(9.5), Inches(0.9), str(uni).upper(), 16, ctx["on_dark"], bold=True, font=ctx["font"])
    if kaf and kaf != uni:
        _add_text(s, Inches(0.9), Inches(1.55), Inches(9.5), Inches(0.6), str(kaf), 13, ctx["accent2"], font=ctx["font"])
    _add_text(s, Inches(0.9), Inches(2.7), Inches(11.0), Inches(2.2), d.get("topic", "Mavzu"), 40, ctx["on_dark"], bold=True, anchor=MSO_ANCHOR.MIDDLE, font=ctx["font"])
    _add_rect(s, Inches(0.95), Inches(4.95), Inches(2.6), Inches(0.08), ctx["accent"])
    _title_info(s, ctx, d, ctx["on_dark"])
    return s


def _title_split(prs, ctx, d):
    s = _blank_slide(prs)
    _fill_background(s, WHITE)
    _add_rect(s, Inches(0), Inches(0), Inches(6.0), SLIDE_H, ctx["primary"])
    _add_rect(s, Inches(6.0), Inches(0), Inches(0.12), SLIDE_H, ctx["accent"])
    uni = d.get("university_name") or d.get("uni_kafedra", "")
    _add_text(s, Inches(0.7), Inches(0.8), Inches(4.8), Inches(0.9), str(uni).upper(), 14, ctx["on_dark"], bold=True, font=ctx["font"])
    _add_text(s, Inches(0.7), Inches(2.5), Inches(4.8), Inches(2.6), d.get("topic", "Mavzu"), 34, ctx["on_dark"], bold=True, anchor=MSO_ANCHOR.MIDDLE, font=ctx["font"])
    _add_rect(s, Inches(0.72), Inches(5.2), Inches(1.8), Inches(0.1), ctx["accent"])
    info = [f"Bajardi:  {d.get('student_fio', '')}", f"Guruh:  {d.get('student_group', '')}", f"Yil:  {d.get('year', '')}"]
    tb = s.shapes.add_textbox(Inches(6.6), Inches(2.8), Inches(6.0), Inches(2.0))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(info):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run()
        r.text = line
        r.font.size = Pt(16)
        r.font.name = ctx["font"]
        r.font.color.rgb = ctx["text"]
        p.space_after = Pt(4)
    return s


def _title_band(prs, ctx, d):
    s = _blank_slide(prs)
    _fill_background(s, WHITE)
    _add_rect(s, Inches(0), Inches(0), SLIDE_W, Inches(2.0), ctx["primary"])
    _add_rect(s, Inches(0), Inches(2.0), SLIDE_W, Inches(0.1), ctx["accent"])
    uni = d.get("university_name") or d.get("uni_kafedra", "")
    _add_text(s, Inches(0.9), Inches(0.6), Inches(11.0), Inches(0.9), str(uni).upper(), 16, ctx["on_dark"], bold=True, anchor=MSO_ANCHOR.MIDDLE, font=ctx["font"])
    _add_text(s, Inches(0.9), Inches(2.9), Inches(11.5), Inches(2.0), d.get("topic", "Mavzu"), 38, ctx["primary"], bold=True, anchor=MSO_ANCHOR.MIDDLE, font=ctx["font"])
    _add_rect(s, Inches(0.95), Inches(5.0), Inches(2.6), Inches(0.08), ctx["accent"])
    _title_info(s, ctx, d, ctx["text"])
    return s


def _title_light(prs, ctx, d):
    s = _blank_slide(prs)
    _fill_background(s, ctx["bg"] if ctx["bg"] != ctx["primary"] else WHITE)
    _add_rect(s, Inches(0), Inches(0), Inches(0.28), SLIDE_H, ctx["accent"])
    uni = d.get("university_name") or d.get("uni_kafedra", "")
    _add_text(s, Inches(1.0), Inches(0.9), Inches(10.0), Inches(0.9), str(uni).upper(), 15, ctx["primary"], bold=True, font=ctx["font"])
    _add_text(s, Inches(1.0), Inches(2.7), Inches(11.0), Inches(2.0), d.get("topic", "Mavzu"), 38, ctx["primary"], bold=True, anchor=MSO_ANCHOR.MIDDLE, font=ctx["font"])
    _add_rect(s, Inches(1.05), Inches(4.8), Inches(2.4), Inches(0.07), ctx["accent"])
    _title_info(s, ctx, d, ctx["text"])
    return s


def _title_geo(prs, ctx, d):
    """Zamonaviy — o'ng tomonda katta rangli blok + geometrik aksent."""
    s = _blank_slide(prs)
    _fill_background(s, WHITE)
    _add_rect(s, Inches(8.5), Inches(0), Inches(4.83), SLIDE_H, ctx["primary"])
    _add_rect(s, Inches(7.4), Inches(2.3), Inches(1.5), Inches(1.5), ctx["accent"], shape=MSO_SHAPE.OVAL)
    _add_rect(s, Inches(9.4), Inches(5.2), Inches(1.1), Inches(1.1), ctx["accent2"], shape=MSO_SHAPE.OVAL)
    uni = d.get("university_name") or d.get("uni_kafedra", "")
    _add_text(s, Inches(0.9), Inches(0.9), Inches(7.0), Inches(0.9), str(uni).upper(), 15, ctx["primary"], bold=True, font=ctx["font"])
    _add_text(s, Inches(0.9), Inches(3.0), Inches(7.0), Inches(2.2), d.get("topic", "Mavzu"), 36, ctx["primary"], bold=True, anchor=MSO_ANCHOR.MIDDLE, font=ctx["font"])
    _add_rect(s, Inches(0.95), Inches(5.2), Inches(2.4), Inches(0.08), ctx["accent"])
    _title_info(s, ctx, d, ctx["text"])
    return s


def _title_centerdark(prs, ctx, d):
    """Tungi — to'liq quyuq fon, markazlashgan katta sarlavha."""
    s = _blank_slide(prs)
    _fill_background(s, ctx["primary"])
    uni = d.get("university_name") or d.get("uni_kafedra", "")
    _add_text(s, Inches(1.0), Inches(1.1), Inches(11.3), Inches(0.7), str(uni).upper(), 15, ctx["accent2"], bold=True, align=PP_ALIGN.CENTER, font=ctx["font"])
    _add_text(s, Inches(1.0), Inches(2.6), Inches(11.3), Inches(2.0), d.get("topic", "Mavzu"), 42, ctx["on_dark"], bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=ctx["font"])
    _add_rect(s, Inches(5.86), Inches(4.7), Inches(1.6), Inches(0.08), ctx["accent"])
    _add_text(s, Inches(1.0), Inches(5.4), Inches(11.3), Inches(1.2),
              f"{d.get('student_fio', '')}  •  {d.get('student_group', '')}  •  {d.get('year', '')}",
              15, ctx["on_dark"], align=PP_ALIGN.CENTER, font=ctx["font"])
    return s


def _title_diagonal(prs, ctx, d):
    """Ijodiy — assimetrik uchburchak shakllar."""
    s = _blank_slide(prs)
    _fill_background(s, WHITE)
    _add_rect(s, Inches(0), Inches(0), Inches(3.4), Inches(3.4), ctx["accent"], shape=MSO_SHAPE.RIGHT_TRIANGLE)
    tri = _add_rect(s, Inches(9.9), Inches(4.1), Inches(3.43), Inches(3.4), ctx["accent2"], shape=MSO_SHAPE.RIGHT_TRIANGLE)
    tri.rotation = 180
    _add_rect(s, Inches(11.0), Inches(0.6), Inches(1.1), Inches(1.1), ctx["primary"], shape=MSO_SHAPE.OVAL)
    uni = d.get("university_name") or d.get("uni_kafedra", "")
    _add_text(s, Inches(1.4), Inches(1.2), Inches(9.0), Inches(0.8), str(uni).upper(), 15, ctx["primary"], bold=True, font=ctx["font"])
    _add_text(s, Inches(1.4), Inches(2.7), Inches(9.5), Inches(2.0), d.get("topic", "Mavzu"), 38, ctx["primary"], bold=True, anchor=MSO_ANCHOR.MIDDLE, font=ctx["font"])
    _add_rect(s, Inches(1.45), Inches(4.8), Inches(2.4), Inches(0.08), ctx["accent"])
    _title_info(s, ctx, d, ctx["text"])
    return s


def _title_elegant(prs, ctx, d):
    """Nafis — krem fon, markazlashgan serif, ikki nozik chiziq."""
    s = _blank_slide(prs)
    _fill_background(s, ctx["bg"] if ctx["bg"] != ctx["primary"] else CREAM)
    uni = d.get("university_name") or d.get("uni_kafedra", "")
    _add_text(s, Inches(1.0), Inches(1.4), Inches(11.3), Inches(0.7), str(uni).upper(), 14, ctx["primary"], align=PP_ALIGN.CENTER, font=ctx["font"])
    _add_rect(s, Inches(5.16), Inches(2.7), Inches(3.0), Inches(0.03), ctx["accent"])
    _add_text(s, Inches(1.0), Inches(2.9), Inches(11.3), Inches(1.6), d.get("topic", "Mavzu"), 36, ctx["primary"], bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=ctx["font"], italic=True)
    _add_rect(s, Inches(5.16), Inches(4.6), Inches(3.0), Inches(0.03), ctx["accent"])
    _add_text(s, Inches(1.0), Inches(5.2), Inches(11.3), Inches(1.0),
              f"{d.get('student_fio', '')}  •  {d.get('student_group', '')}  •  {d.get('year', '')}",
              14, ctx["text"], align=PP_ALIGN.CENTER, font=ctx["font"])
    return s


def _title_stat(prs, ctx, d):
    """Infografik — sarlavha + pastda 3 ta dekorativ blok."""
    s = _blank_slide(prs)
    _fill_background(s, WHITE)
    _add_rect(s, Inches(0.9), Inches(0.85), Inches(0.18), Inches(0.7), ctx["accent"])
    uni = d.get("university_name") or d.get("uni_kafedra", "")
    _add_text(s, Inches(1.25), Inches(0.85), Inches(10.5), Inches(0.7), str(uni).upper(), 15, ctx["primary"], bold=True, anchor=MSO_ANCHOR.MIDDLE, font=ctx["font"])
    _add_text(s, Inches(0.9), Inches(1.9), Inches(11.4), Inches(1.6), d.get("topic", "Mavzu"), 36, ctx["primary"], bold=True, anchor=MSO_ANCHOR.MIDDLE, font=ctx["font"])
    cols = [ctx["accent"], ctx["accent2"], ctx["primary"]]
    for i in range(3):
        x = Inches(0.9 + i * 3.95)
        _add_rect(s, x, Inches(3.9), Inches(3.7), Inches(1.5), cols[i], shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    _add_text(s, Inches(0.9), Inches(5.7), Inches(11.4), Inches(0.8),
              f"{d.get('student_fio', '')}  •  {d.get('student_group', '')}  •  {d.get('year', '')}",
              14, ctx["text"], font=ctx["font"])
    return s


def _title_frame(prs, ctx, d):
    """Rasm asosida — quyuq fon + qalin aksent ramka (poster uslubi)."""
    s = _blank_slide(prs)
    _fill_background(s, ctx["primary"])
    bw = Inches(0.22)
    _add_rect(s, Inches(0.5), Inches(0.45), Inches(12.33), bw, ctx["accent"])
    _add_rect(s, Inches(0.5), Inches(6.83), Inches(12.33), bw, ctx["accent"])
    _add_rect(s, Inches(0.5), Inches(0.45), bw, Inches(6.6), ctx["accent"])
    _add_rect(s, Inches(12.61), Inches(0.45), bw, Inches(6.6), ctx["accent"])
    uni = d.get("university_name") or d.get("uni_kafedra", "")
    _add_text(s, Inches(1.2), Inches(1.3), Inches(10.9), Inches(0.8), str(uni).upper(), 15, ctx["accent2"], bold=True, align=PP_ALIGN.CENTER, font=ctx["font"])
    _add_text(s, Inches(1.2), Inches(2.8), Inches(10.9), Inches(1.8), d.get("topic", "Mavzu"), 40, ctx["on_dark"], bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=ctx["font"])
    _add_text(s, Inches(1.2), Inches(5.4), Inches(10.9), Inches(0.9),
              f"{d.get('student_fio', '')}  •  {d.get('student_group', '')}  •  {d.get('year', '')}",
              14, ctx["on_dark"], align=PP_ALIGN.CENTER, font=ctx["font"])
    return s


_TITLES = {
    "block": _title_block, "minimal": _title_light, "split": _title_split,
    "band": _title_band, "geo": _title_geo, "centerdark": _title_centerdark,
    "diagonal": _title_diagonal, "elegant": _title_elegant, "stat": _title_stat,
    "frame": _title_frame,
}


# ---------------------------------------------------------------------------
# Kontent / grafik / yakun slaydlari
# ---------------------------------------------------------------------------
def _build_content_slide(prs, ctx, title, bullets, index, total, topic, image_path, icons, parity, notes=None):
    s = _blank_slide(prs)
    _fill_background(s, ctx["bg"])
    if not bullets:
        bullets = ["Ma'lumot generatsiya qilinmadi."]
    hi = bool(image_path) and os.path.exists(image_path or "")
    fam = ctx["family"]
    try:
        _FAMILIES[fam](s, ctx, title, bullets, hi, image_path, icons, parity)
    except Exception as e:
        print(f"⚠️ '{fam}' oila xato, standartga qaytildi: {e}")
        _fam_standard(s, ctx, title, bullets, hi, image_path, icons, parity)
    _add_footer(s, ctx, index, total, topic, show_topic=(fam not in _PANEL_FAMILIES))
    if notes:
        try:
            s.notes_slide.notes_text_frame.text = notes
        except Exception:
            pass
    return s


def _build_chart_slide(prs, ctx, cd, index, total, topic, chart_type="column"):
    s = _blank_slide(prs)
    _fill_background(s, ctx["bg"])
    theme = ctx["theme"]
    _header(s, ctx, cd.get("title") or "Statistik ko'rsatkichlar")
    xl = CHART_TYPE_MAP.get(chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED)
    chd = CategoryChartData()
    chd.categories = cd["categories"]
    if xl == XL_CHART_TYPE.PIE:
        f = cd["series"][0]
        chd.add_series(f["name"], f["values"])
    else:
        for ser in cd["series"]:
            chd.add_series(ser["name"], ser["values"])
    gf = s.shapes.add_chart(xl, Inches(1.2), Inches(1.95), Inches(10.9), Inches(4.8), chd)
    ch = gf.chart
    ch.has_title = False
    if xl == XL_CHART_TYPE.PIE:
        ch.has_legend = True
        ch.legend.position = XL_LEGEND_POSITION.RIGHT
        ch.legend.include_in_layout = False
        try:
            pts = ch.plots[0].series[0].points
            pal = [theme["accent"], theme["accent2"], theme["primary"], theme["muted"], _C(0x93, 0xC5, 0xFD)]
            for i, pt in enumerate(pts):
                pt.format.fill.solid()
                pt.format.fill.fore_color.rgb = pal[i % len(pal)]
        except Exception:
            pass
    else:
        ch.has_legend = len(cd["series"]) > 1
        if ch.has_legend:
            ch.legend.position = XL_LEGEND_POSITION.BOTTOM
            ch.legend.include_in_layout = False
        try:
            pal = [theme["accent"], theme["accent2"], theme["primary"]]
            for si, ps in enumerate(ch.series):
                ps.format.fill.solid()
                ps.format.fill.fore_color.rgb = pal[si % len(pal)]
        except Exception:
            pass
    _add_footer(s, ctx, index, total, topic)
    return s


def _build_closing_slide(prs, ctx, d):
    s = _blank_slide(prs)
    _fill_background(s, ctx["primary"])
    _add_rect(s, Inches(0), Inches(0), Inches(0.28), SLIDE_H, ctx["accent"])
    _add_text(s, Inches(0.9), Inches(2.7), Inches(11.5), Inches(1.6), "E'tiboringiz uchun rahmat!", 44, ctx["on_dark"], bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=ctx["font"])
    _add_rect(s, Inches(5.86), Inches(4.4), Inches(1.6), Inches(0.08), ctx["accent"])
    _add_text(s, Inches(0.9), Inches(4.7), Inches(11.5), Inches(0.8), str(d.get("student_fio", "")), 18, ctx["accent2"], align=PP_ALIGN.CENTER, font=ctx["font"])
    return s


# ---- Qo'shimcha professional slaydlar ----
def _build_agenda_slide(prs, ctx, titles, index, total, topic):
    s = _blank_slide(prs)
    _fill_background(s, ctx["bg"])
    _header(s, ctx, "Reja")
    items = [t for t in titles][:8]
    box = s.shapes.add_textbox(Inches(1.0), Inches(1.95), Inches(11.3), Inches(4.8))
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    for i, t in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(10)
        m = p.add_run()
        m.text = f"{i+1}.  "
        m.font.size = Pt(18)
        m.font.bold = True
        m.font.name = ctx["font"]
        m.font.color.rgb = ctx["accent"]
        r = p.add_run()
        r.text = _shorten(t, 90)
        r.font.size = Pt(18)
        r.font.name = ctx["font"]
        r.font.color.rgb = ctx["text"]
    _add_footer(s, ctx, index, total, topic)
    return s


def _build_section_divider(prs, ctx, label, subtitle):
    s = _blank_slide(prs)
    _fill_background(s, ctx["primary"])
    _add_rect(s, Inches(0), Inches(0), Inches(0.28), SLIDE_H, ctx["accent"])
    _add_text(s, Inches(1.2), Inches(2.4), Inches(10.9), Inches(1.2), label, 26, ctx["accent2"], bold=True, anchor=MSO_ANCHOR.MIDDLE, font=ctx["font"])
    _add_rect(s, Inches(1.25), Inches(3.7), Inches(2.2), Inches(0.1), ctx["accent"])
    _add_text(s, Inches(1.2), Inches(3.95), Inches(10.9), Inches(1.8), _shorten(subtitle, 80), 34, ctx["on_dark"], bold=True, anchor=MSO_ANCHOR.TOP, font=ctx["font"])
    return s


def _build_references_slide(prs, ctx, refs_text, index, total, topic):
    s = _blank_slide(prs)
    _fill_background(s, ctx["bg"])
    _header(s, ctx, "Foydalanilgan adabiyotlar")
    refs = [r.strip() for r in (refs_text or "").split("\n") if r.strip()][:8]
    if not refs:
        refs = ["Internet manbalari va darsliklar."]
    box = s.shapes.add_textbox(Inches(1.0), Inches(1.95), Inches(11.3), Inches(4.8))
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    for i, r in enumerate(refs):
        r = r.lstrip("0123456789.)-•* ").strip()
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(8)
        p.line_spacing = 1.05
        m = p.add_run()
        m.text = f"{i+1}.  "
        m.font.size = Pt(14)
        m.font.bold = True
        m.font.name = ctx["font"]
        m.font.color.rgb = ctx["accent"]
        rn = p.add_run()
        rn.text = _shorten(r, 150)
        rn.font.size = Pt(14)
        rn.font.name = ctx["font"]
        rn.font.color.rgb = ctx["text"]
    _add_footer(s, ctx, index, total, topic)
    return s


def _build_qa_slide(prs, ctx):
    s = _blank_slide(prs)
    _fill_background(s, ctx["primary"])
    _add_rect(s, Inches(0), Inches(0), Inches(0.28), SLIDE_H, ctx["accent"])
    _add_text(s, Inches(0.9), Inches(2.6), Inches(11.5), Inches(1.6), "Savollaringiz bormi?", 44, ctx["on_dark"], bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=ctx["font"])
    _add_rect(s, Inches(5.86), Inches(4.3), Inches(1.6), Inches(0.08), ctx["accent"])
    _add_text(s, Inches(0.9), Inches(4.6), Inches(11.5), Inches(0.7), "Javob berishdan mamnun bo'laman", 18, ctx["accent2"], align=PP_ALIGN.CENTER, font=ctx["font"])
    return s


def _build_table_slide(prs, ctx, td, index, total, topic):
    s = _blank_slide(prs)
    _fill_background(s, ctx["bg"])
    _header(s, ctx, td.get("title") or "Taqqoslash jadvali")
    headers = td.get("headers") or []
    rows = td.get("rows") or []
    rows = [r for r in rows if r][:5]
    ncol = max(1, len(headers))
    nrow = len(rows) + 1
    gx = s.shapes.add_table(nrow, ncol, Inches(1.0), Inches(2.0), Inches(11.3), Inches(0.8 + 0.7 * len(rows)))
    table = gx.table
    for c in range(ncol):
        cell = table.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = ctx["primary"]
        cell.text = str(headers[c]) if c < len(headers) else ""
        para = cell.text_frame.paragraphs[0]
        para.font.size = Pt(14)
        para.font.bold = True
        para.font.name = ctx["font"]
        para.font.color.rgb = ctx["on_dark"]
    for ri, row in enumerate(rows, start=1):
        for c in range(ncol):
            cell = table.cell(ri, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = ctx["theme"]["light"] if ri % 2 else WHITE
            cell.text = str(row[c]) if c < len(row) else ""
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(12)
            para.font.name = ctx["font"]
            para.font.color.rgb = ctx["theme"]["text"]
    _add_footer(s, ctx, index, total, topic)
    return s


def _build_timeline_slide(prs, ctx, tl, index, total, topic):
    s = _blank_slide(prs)
    _fill_background(s, ctx["bg"])
    _header(s, ctx, tl.get("title") or "Bosqichlar")
    steps = (tl.get("steps") or [])[:5]
    n = max(1, len(steps))
    margin = 1.0
    avail = 11.33
    bw = avail / n
    cy = 4.0
    # bog'lovchi chiziq
    _add_rect(s, Inches(margin + bw / 2), Inches(cy - 0.01), Inches(avail - bw), Inches(0.05), ctx["line"])
    for i, st in enumerate(steps):
        cx = margin + bw * i + bw / 2
        # raqamli doira
        _add_rect(s, Inches(cx - 0.45), Inches(cy - 0.45), Inches(0.9), Inches(0.9), ctx["accent"], shape=MSO_SHAPE.OVAL)
        _add_text(s, Inches(cx - 0.45), Inches(cy - 0.45), Inches(0.9), Inches(0.9), str(i + 1), 20, ctx["on_dark"], bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=ctx["font"])
        # label (tepada)
        _add_text(s, Inches(cx - bw / 2 + 0.1), Inches(2.4), Inches(bw - 0.2), Inches(1.0), _shorten(st.get("label", ""), 30), 14, ctx["title"], bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.BOTTOM, font=ctx["font"])
        # text (pastda)
        _add_text(s, Inches(cx - bw / 2 + 0.1), Inches(4.7), Inches(bw - 0.2), Inches(1.6), _shorten(st.get("text", ""), 70), 12, ctx["text"], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP, font=ctx["font"])
    _add_footer(s, ctx, index, total, topic)
    return s


# ---------------------------------------------------------------------------
# Asosiy funksiya
# ---------------------------------------------------------------------------
async def generate_pptx_file(doc_data, presentation_content, temp_dir,
                             theme_name=None, template_name=None, options=None, theme_path=None):
    """
    doc_data: {topic, student_fio, student_group, uni_kafedra, university_name, year}
    presentation_content: [{"title", "content", "image"}]
    theme_name: 15 rangdan biri (None -> tasodifiy)
    template_name: 10 shablondan biri (None -> 'classic')
    options: {"images","icons","chart_type","chart_count"}
    """
    options = options or {}
    use_icons = bool(options.get("icons"))
    chart_type = options.get("chart_type")
    chart_count = int(options.get("chart_count") or 0)
    if not chart_type:
        chart_count = 0

    if theme_name not in THEMES:
        theme_name = random.choice(list(THEMES.keys()))
    if template_name not in TEMPLATES:
        template_name = "classic"
    theme = THEMES[theme_name]
    tpl = TEMPLATES[template_name]
    ctx = _ctx(theme, tpl)

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    topic = doc_data.get("topic", "Mavzu")

    # Titul slayd (shablonga mos variant)
    _TITLES.get(ctx["title_kind"], _title_block)(prs, ctx, doc_data)

    # Grafik ma'lumotlari
    chart_dicts = []
    if chart_count > 0:
        try:
            gemini = GeminiService()
            for _ in range(chart_count):
                data = await gemini.generate_chart_data(topic)
                if data and data.get("categories") and data.get("series"):
                    chart_dicts.append(data)
        except Exception as e:
            print(f"⚠️ Grafik ma'lumoti olinmadi: {e}")
            chart_dicts = []

    # Qo'shimcha professional opsiyalar
    use_structure = bool(options.get("structure"))
    use_refs_qa = bool(options.get("refs_qa"))
    references_text = options.get("references_text") or ""
    table_data = options.get("table_data") if options.get("visuals") else None
    timeline_data = options.get("timeline_data") if options.get("visuals") else None

    tc = len(presentation_content)
    titles = [it.get("title", f"Slayd {i+1}") for i, it in enumerate(presentation_content)]

    # Grafik joylashuvi
    chart_positions = {}
    if chart_dicts and tc > 0:
        step = max(1, tc // (len(chart_dicts) + 1))
        for ci in range(len(chart_dicts)):
            pos = min(tc - 1, step * (ci + 1))
            while pos in chart_positions and pos < tc - 1:
                pos += 1
            chart_positions[pos] = ci

    # Bo'lim ajratuvchilar (boshlang'ich slaydni hisobga olmagan holda)
    divider_positions = {}
    if use_structure and tc >= 6:
        b1, b2 = tc // 3, (2 * tc) // 3
        if b1 > 0:
            divider_positions[b1] = ("II BO'LIM", titles[b1] if b1 < tc else "")
        if b2 > b1:
            divider_positions[b2] = ("III BO'LIM", titles[b2] if b2 < tc else "")

    table_pos = (tc // 2) if table_data else -1
    timeline_pos = ((2 * tc) // 3) if timeline_data else -1

    # Umumiy slaydlar soni (footer uchun)
    total_slides = (1
                    + (1 if use_structure else 0)
                    + len(divider_positions)
                    + tc
                    + len(chart_dicts)
                    + (1 if table_data else 0)
                    + (1 if timeline_data else 0)
                    + (1 if use_refs_qa else 0)   # adabiyotlar
                    + (1 if use_refs_qa else 0)   # Q&A
                    + 1)                          # yakun
    counter = 1  # titul

    # Reja (mundarija)
    if use_structure:
        counter += 1
        _build_agenda_slide(prs, ctx, titles, counter, total_slides, topic)

    used = set()
    for i, item in enumerate(presentation_content):
        if i in divider_positions:
            label, sub = divider_positions[i]
            counter += 1
            _build_section_divider(prs, ctx, label, sub)
        if i in chart_positions:
            ci = chart_positions[i]
            if ci not in used:
                counter += 1
                _build_chart_slide(prs, ctx, chart_dicts[ci], counter, total_slides, topic, chart_type=chart_type)
                used.add(ci)
        counter += 1
        bullets = _clean_bullets(item.get("content", ""))
        _build_content_slide(prs, ctx, item.get("title", f"Slayd {i+1}"), bullets,
                             counter, total_slides, topic, item.get("image"), use_icons,
                             parity=i, notes=item.get("notes"))
        if i == table_pos and table_data:
            counter += 1
            _build_table_slide(prs, ctx, table_data, counter, total_slides, topic)
        if i == timeline_pos and timeline_data:
            counter += 1
            _build_timeline_slide(prs, ctx, timeline_data, counter, total_slides, topic)

    for ci in range(len(chart_dicts)):
        if ci not in used:
            counter += 1
            _build_chart_slide(prs, ctx, chart_dicts[ci], counter, total_slides, topic, chart_type=chart_type)

    # Adabiyotlar + Q&A
    if use_refs_qa:
        counter += 1
        _build_references_slide(prs, ctx, references_text, counter, total_slides, topic)
        counter += 1
        _build_qa_slide(prs, ctx)

    _build_closing_slide(prs, ctx, doc_data)

    os.makedirs(temp_dir, exist_ok=True)
    fp = os.path.join(temp_dir, f"PPTX_{uuid.uuid4().hex[:8]}.pptx")
    prs.save(fp)
    print(f"✅ Tayyor — shablon: {tpl['name']}, rang: {theme['name']} -> {fp}")
    return fp
